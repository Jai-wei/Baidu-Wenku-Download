import os
import re
import json
import time
import requests
from lxml import etree
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import RGBColor
from docx.shared import Pt
from docx.shared import Twips
from io import BytesIO
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
import requests
import io
import PIL
from PIL import Image, ImageOps, ImageChops
import uuid
from docx.shared import Cm
from pptx import Presentation
from pptx.util import Inches


class BDWK(object):
    __qualname__ = 'BDWK'

    def __init__(self, url):
        self.title = None
        self.url = url
        self.docType = None
        self.headers = {'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_12_4) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/62.0.3202.75 Safari/537.36'}
        self.get_response_content(self.url)
        self.type_and_title()

    def get_response_content(self, url):
        try :
            html = requests.get(url, headers=(self.headers))
            return html.content
        except Exception as e:
            print(e)

    def type_and_title(self):
        source_html = self.get_response_content(self.url)
        content = source_html.decode('gbk')
        self.docType = re.findall("docType.*?\\:.*?\\'(.*?)\\'\\,", content)[0]
        self.title = re.findall("title.*?\\:.*?\\'(.*?)\\'\\,", content)[0]

class WKTXT(BDWK):
    __qualname__ = 'BDWKTXT'

    def __init__(self = None, url = None):
        super().__init__(url)
        self.doc_id = None
        self.get_txt(url)

    def get_txt(self,url):

        self.doc_id = re.findall('view/(.*).html', url)[0]
        token_url = 'https://wenku.baidu.com/api/doc/getdocinfo?callback=cb&doc_id=' + self.doc_id
        x = self.get_response_content(token_url).decode()
        first_json = json.loads(re.match('.*?\\((\\{.*?\\})\\).*', self.get_response_content(token_url).decode()).group(1))
        md5sum = first_json['md5sum']
        pn = first_json['docInfo']['totalPageNum']
        rsign = first_json['rsign']
        target_url = 'https://wkretype.bdimg.com/retype/text/'+ self.doc_id +'?rn='+ pn +'&type=txt'+ md5sum +'&rsign='+rsign
        txt = requests.get(target_url).text     # yes
        jsons = json.loads(txt)     # yes
        texts = []
        for i in range(len(jsons)):
            #print(str(jsons[i]['parags']))
            text = re.findall("'c': ['\"](.*?)['\"],", str(jsons[i]))   # https://regex101.com/
            #print(text)
            texts.extend(text)
        filename = './download/' + self.title + '.txt'
        with open(filename, 'a', encoding='utf-8') as f:
            for i in range(0, len(texts)):
                #print(i)
                #print('\n')
                texts[i] = texts[i].replace('\\r', '\r')
                texts[i] = texts[i].replace('\\n', '\n')
                #print(texts[i])
                f.write(texts[i])
        print("文档保存在" + filename)

class WKDOC(BDWK):
    def __init__(self, url):
        super().__init__(url)
        self.pure_addr_list = list()
        self.pure_addr_list_pic = list()
        self.space_lign = 4
        self.space_lign_high = 2
        self.line_list = list()

    def color(self, value):
        digit = list(map(str, range(10))) + list('ABCDEF')
        if isinstance(value, tuple):
            string = '#'
            for i in value:
                a1 = i // 16
                a2 = i % 16
                string += digit[a1] + digit[a2]

            return string
        if isinstance(value, str):
            value = value.upper()
            a1 = digit.index(value[1]) * 16 + digit.index(value[2])
            a2 = digit.index(value[3]) * 16 + digit.index(value[4])
            a3 = digit.index(value[5]) * 16 + digit.index(value[6])
            return RGBColor(a1, a2, a3)

    def get_pure_addr_list(self):
        source_html = self.get_response_content(self.url).decode('gbk')
        all_addr = re.findall('wkbjcloudbos\\.bdimg\\.com.*?json.*?Expire.*?\\}', source_html)
        pure_addr_list_temp = list()
        self.title = etree.HTML(source_html).xpath('//title/text()')[0].strip()
        for addr in all_addr:
            addr = 'https://' + addr.replace('\\\\\\/', '/')
            addr = addr[:-5]
            pure_addr_list_temp.append(addr)
        self.pure_addr_list = list(set(pure_addr_list_temp))
        self.pure_addr_list.sort(key=(pure_addr_list_temp.index))

        str = source_html[source_html.find('WkInfo.htmlUrls = ') + len('WkInfo.htmlUrls = '):]
        str_end = str[str.find('png'):str.find("}]}'")]
        all_addr = re.findall('wkbjcloudbos\\.bdimg\\.com.*?png.*?Expire.*?\\}', str_end)
        pure_addr_list_img = list()
        self.title = etree.HTML(source_html).xpath('//title/text()')[0].strip()
        for addr in all_addr:
            addr = 'https://' + addr.replace('\\\\\\/', '/')
            addr = addr[:-5]
            pure_addr_list_img.append(addr)

        self.pure_addr_list_pic = list(set(pure_addr_list_img))
        self.pure_addr_list_pic.sort(key=(pure_addr_list_img.index))
        return (self.pure_addr_list, self.pure_addr_list_pic)

    def sort_by_y(self, first):
        temp_return_y = 0
        temp_return_h = 0
        for line in self.line_list:
            if first['p']['y'] + first['p']['h'] / 2 < line[0] + line[1] / 2:
                if first['p']['y'] + first['p']['h'] / 2 + self.space_lign_high > line[0] + line[1] / 2:
                    temp_return_y = line[0]
                    temp_return_h = line[1]
            if first['p']['y'] + first['p']['h'] / 2 > line[0] + line[1] / 2 and first['p']['y'] + first['p'][
                'h'] / 2 - self.space_lign_high < line[0] + line[1] / 2:
                temp_return_y = line[0]
                temp_return_h = line[1]

        if temp_return_h == 0:
            temp_return_y = first['p']['y']
            temp_return_h = first['p']['h']
            self.line_list.append([first['p']['y'], first['p']['h']])
        return temp_return_y * 10000 + first['p']['x']

    def trim(im):
        bg = Image.new(im.mode, im.size, im.getpixel((0, 0)))
        diff = ImageChops.difference(im, bg)
        diff = ImageChops.add(diff, diff, 2.0, -100)
        bbox = diff.getbbox()
        if bbox:
            return im.crop(bbox)

    def sort_by_x(self, first):
        return first['p']['x']

    def get_json_content(self, url_list, url_list_pic, doc_name='test.doc', doc_path=''):
        content = ''
        result = ''
        sum = len(url_list)
        i = 1
        content_array = []
        Img_temp = None
        for pure_addr in url_list:
            print('正在下载第%d-%d条数据' % (i, sum - i))
            i += 1
            Img_temp = None
            self.line_list = list()
            try:
                #print(repr(pure_addr))
                content = self.get_response_content(pure_addr).decode()
                content = re.match('.*?\\((.*)\\)$', content).group(1)
                style_array = []
                all_font_info = json.loads(content)['font']
                all_style_info = json.loads(content)['style']
                for style_info in all_style_info:
                    try:
                        dic_style = {}
                        style_info_index = style_info['c'][(-1)]
                        for style_info_c_item in all_style_info:
                            if style_info_index in style_info_c_item['c']:
                                dic_style.update(style_info_c_item['s'])

                        style_array.append(dic_style)
                    except Exception as e:
                        print(e)

                for style_item in style_array:
                    if style_item.get('font-family'):
                        value_font = all_font_info[style_item.get('font-family')]
                        style_item['font-family'] = value_font

                all_body_info = json.loads(content)['body']
                x_point = 0
                string_temp = None
                all_body_info.sort(key=(self.sort_by_y))
                for index_body in range(len(all_body_info)):
                    body_info = all_body_info[index_body]
                    try:
                        if isinstance(body_info['c'], dict):        # 检查图片
                            if not Img_temp:
                                test = self.get_response_content(url_list_pic[(i - 2)])     # 此处数组越界,不影响
                                f = io.BytesIO(test)
                                Img_temp = Image.open(f)
                            if body_info['t'] != 'pic' or body_info['s'] == None:
                                continue
                            region = (
                                body_info['c']['ix'], body_info['c']['iy'],
                                body_info['c']['iw'] + body_info['c']['ix'] - 1,
                                body_info['c']['ih'] + body_info['c']['iy'] - 1)
                            cropImg = Img_temp.crop(region)         # 图片裁剪
                            image = cropImg.convert('RGB')
                            ivt_image = ImageOps.invert(image)
                            bbox = image.getbbox()                  # 检测图像边界
                            left = bbox[0]
                            top = bbox[1]
                            right = bbox[2]
                            bottom = bbox[3]
                            cropImg = cropImg.crop([left, top, right, bottom])
                            output_buffer = BytesIO()
                            width, crop_high = cropImg.size
                            pic_temp = io.BytesIO()
                            cropImg.save(pic_temp, format='PNG')
                            content_array.append([pic_temp, 'pic', width / 39])
                            continue
                        style_temp = {}
                        for style_item in body_info['r']:               # r代表字体编号
                            style_temp.update(style_array[style_item])
                        if body_info.get('s'):                          # s字体样式
                            style_temp.update(body_info.get('s'))

                        if isinstance(body_info['c'], str):             # 检查汉字
                            if body_info['t'] == 'word':
                                content_array.append([body_info['c'], 'str', style_temp])
                            if body_info['t'] != 'word':
                                continue

                        if body_info.get('ps'):
                            if body_info['ps'].get('_enter'):
                                if body_info['ps'].get('_enter') == 1:
                                    content_array.append([body_info['c'], style_temp])
                                    if index_body < len(all_body_info) - 1:
                                        if body_info['p']['x'] + body_info['p']['w'] - self.space_lign > \
                                                all_body_info[(index_body + 1)]['p']['x']:
                                            content_array.append([None, None])
                                            x_point = 0
                                        if body_info['p']['x'] < x_point:
                                            if string_temp == ' ':
                                                content_array.append([None, None])
                                        x_point = body_info['p']['x'] + body_info['p']['w'] - self.space_lign
                                        string_temp = body_info['c']
                                        content_array.append([body_info['c'], style_temp])
                    except Exception as e:
                        print(e)

            except Exception as e:
                print(e)

        document = Document()
        p = document.add_paragraph()
        for content_array_item in content_array:
            if content_array_item[0] == None:           # ps = _enter,换行
                p = document.add_paragraph()
                continue
            if isinstance(content_array_item[1], str):
                if content_array_item[1] == 'pic':
                    p.add_run().add_picture((content_array_item[0]), width=(Cm(content_array_item[2])))
                    continue
                if content_array_item[1] == 'str':
                    run = p.add_run(content_array_item[0])
                    if content_array_item[2].get('font-family'):
                        run.font.name = content_array_item[2].get('font-family')
                    if content_array_item[2].get('font-size'):
                        run.font.size = Pt(float(content_array_item[2].get('font-size')) * 0.64)
                    if content_array_item[2].get('color'):
                        run.font.color.rgb = self.color(content_array_item[2].get('color'))
                    run.bold = content_array_item[2].get('bold') == 'true' and True

        if doc_path == '':
            document.save(os.getcwd() + '/download/' + doc_name + '.doc')
        else:
            document.save(doc_path)

class WKPPT(BDWK):

    def __init__(self, url):
        self.all_img_url = list()
        super().__init__(url)

    def get_ppt_json_info(self, doc_name='test.ppt', doc_path=''):
        ppt_source_html = self.get_response_content(self.url)
        content = ppt_source_html.decode('gbk')
        self.docId = re.findall("docId.*?(\\w{24}?)\\'\\,", content)[0]
        source_json_url = 'https://wenku.baidu.com/browse/getbcsurl?doc_id=%s&type=ppt&callback=jimmy' % self.docId
        str_source_json = self.get_response_content(source_json_url).decode()
        pure_str_source_json = re.match('.*?\\((.*?)\\)', str_source_json).group(1)
        source_json = json.loads(pure_str_source_json)
        for j in source_json['list']:
            temp_num_url = list()
            temp_num_url.append(j['zoom'])
            temp_num_url.append(j['page'])
            self.all_img_url.append(temp_num_url)

        pptFile = Presentation()
        for img_url in self.all_img_url:
            slide = pptFile.slides.add_slide(pptFile.slide_layouts[1])
            print('正在下载第%d-%d页' % (img_url[1], len(self.all_img_url) - img_url[1]))
            data = self.get_response_content(img_url[0])
            f = io.BytesIO(data)
            slide.shapes.add_picture(f, Inches(0), Inches(0), Inches(10), Inches(7.5))

        if doc_path == '':
            pptFile.save(os.getcwd() + '/download/' + doc_name + '.ppt')
        else:
            pptFile.save(doc_path)
        print('下载成功!')

def main(url):
    try:
        docType = BDWK(url).docType
    except:
        print("无法解析该网址!")
        os.exit()
    print('已检测到', docType, '文件,进入下载')

    if docType == 'txt':
        txt = WKTXT(url)
        print('您要下载的TXT文件名为:', txt.title)
        txt.get_txt(url)
    elif docType == 'doc':
        word = WKDOC(url)
        print("您要下载的DOC文件名为:", word.title)
        pure_addr_list, pure_addr_list_pic = word.get_pure_addr_list()
        word.get_json_content(pure_addr_list, pure_addr_list_pic, word.title)
    elif docType == 'ppt':
        ppt = WKPPT(url)
        print("您要下载的PPT文件名为:", ppt.title)
        ppt.get_ppt_json_info(ppt.title)
    else:
        print("对不起,暂不支持下载" + docType + '类型的文件')

if __name__ == '__main__':
    space_lign = 4
    print("\n请注意: ppt/word/txt 将会下载到本程序所在根目录下的./download文件夹内!\n")
    print("----------***  2秒后进入下载模式  ***----------\n")
    time.sleep(2)

    if not os.path.exists('./download/'):
        os.makedirs('./download')

    print('\n请输入要下载的摆渡文库文件的链接:')
    url = input()
    main(url)



